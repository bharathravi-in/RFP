import os
import uuid
from flask import Blueprint, request, jsonify, current_app
from flask_jwt_extended import jwt_required, get_jwt_identity
from werkzeug.utils import secure_filename
from ..extensions import db
from ..models import Document, Project, User

bp = Blueprint('documents', __name__)

ALLOWED_EXTENSIONS = {'pdf', 'docx', 'xlsx', 'doc', 'xls', 'ppt', 'pptx'}


@bp.route('', methods=['GET'])
@jwt_required()
def list_documents():
    """List documents for a project."""
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    project_id = request.args.get('project_id')
    if not project_id:
        return jsonify({'error': 'project_id required'}), 400
    
    project = Project.query.get(project_id)
    if not project:
        return jsonify({'error': 'Project not found'}), 404
    
    if project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    documents = Document.query.filter_by(project_id=project_id).order_by(Document.created_at.desc()).all()
    
    return jsonify({
        'documents': [d.to_dict() for d in documents]
    }), 200


def allowed_file(filename):
    """Check if file extension is allowed."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


@bp.route('/upload', methods=['POST'])
@jwt_required()
def upload_document():
    """Upload a document to a project with cloud-agnostic storage."""
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    project_id = request.form.get('project_id')
    
    if not project_id:
        return jsonify({'error': 'Project ID required'}), 400
    
    project = Project.query.get(project_id)
    
    if not project:
        return jsonify({'error': 'Project not found'}), 404
    
    if project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': f'File type not allowed. Allowed: {", ".join(ALLOWED_EXTENSIONS)}'}), 400
    
    original_filename = secure_filename(file.filename)
    ext = original_filename.rsplit('.', 1)[1].lower()
    
    # Use cloud-agnostic storage service
    try:
        from app.services.storage_service import get_storage_service
        import os
        
        storage = get_storage_service()
        
        # Check if using GCP with custom prefix
        if storage.storage_type == 'gcp':
            # Get organization name for folder structure
            from app.models import Organization
            org = Organization.query.get(user.organization_id)
            org_name = org.name.lower().replace(' ', '_').replace('/', '_') if org else f"org_{user.organization_id}"
            
            rfp_req_prefix = os.environ.get('GCP_RFP_REQ_PREFIX', 'rfp_requirement')
            # Structure: rfp_requirement/{org_name}/project_{project_id}/
            org_project_subfolder = f"{org_name}/project_{project_id}"
            
            storage_metadata = storage.provider.upload_with_path(
                file=file,
                original_filename=original_filename,
                prefix=rfp_req_prefix,
                subfolder=org_project_subfolder,
                metadata={
                    'project_id': int(project_id),
                    'uploaded_by': user_id,
                    'organization_id': user.organization_id,
                    'organization_name': org_name
                }
            )
        else:
            # Use default upload for local storage
            storage_metadata = storage.upload(
                file=file,
                original_filename=original_filename,
                metadata={
                    'project_id': int(project_id),
                    'uploaded_by': user_id,
                    'organization_id': user.organization_id
                }
            )
        
        # Create document record with storage info
        document = Document(
            file_id=storage_metadata.file_id,
            filename=storage_metadata.file_name,
            original_filename=original_filename,
            file_url=storage_metadata.file_url,
            storage_type=storage_metadata.storage_type,
            file_type=ext,
            file_size=storage_metadata.file_size,
            content_type=storage_metadata.content_type,
            content_hash=storage_metadata.checksum,
            status='pending',
            embedding_status='pending',
            project_id=int(project_id),
            uploaded_by=user_id,
            file_metadata={
                'storage': storage_metadata.to_dict()
            }
        )
        
    except Exception as storage_error:
        # Fallback to legacy DB storage if storage service fails
        current_app.logger.warning(f"Storage service failed, using DB storage: {storage_error}")
        
        file.seek(0)  # Reset file position
        file_data = file.read()
        file_size = len(file_data)
        unique_filename = f"{uuid.uuid4().hex}.{ext}"
        
        document = Document(
            file_id=str(uuid.uuid4()),
            filename=unique_filename,
            original_filename=original_filename,
            file_data=file_data,  # Store binary content in DB
            storage_type='database',
            file_type=ext,
            file_size=file_size,
            status='pending',
            embedding_status='pending',
            project_id=int(project_id),
            uploaded_by=user_id
        )
    
    db.session.add(document)
    db.session.commit()
    
    # Auto-trigger document parsing
    parse_result = None
    try:
        from .documents import _parse_document_internal
        parse_result = _parse_document_internal(document)
    except Exception as e:
        parse_result = {'error': str(e)}
    
    # Trigger background embedding task
    try:
        from app.tasks import process_document_embeddings_task
        process_document_embeddings_task.delay(document.id, user.organization_id)
        embedding_triggered = True
    except Exception as e:
        current_app.logger.warning(f"Failed to trigger embedding task: {e}")
        embedding_triggered = False
    
    return jsonify({
        'message': 'Document uploaded and processing started',
        'document': document.to_dict(),
        'parse_result': parse_result,
        'embedding_triggered': embedding_triggered
    }), 201


@bp.route('/<int:document_id>', methods=['GET'])
@jwt_required()
def get_document(document_id):
    """Get document details."""
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    document = Document.query.get(document_id)
    
    if not document:
        return jsonify({'error': 'Document not found'}), 404
    
    if document.project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    return jsonify({
        'document': document.to_dict()
    }), 200


@bp.route('/<int:document_id>/parse', methods=['POST'])
@jwt_required()
def parse_document(document_id):
    """Trigger document parsing and question extraction."""
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    document = Document.query.get(document_id)
    
    if not document:
        return jsonify({'error': 'Document not found'}), 404
    
    if document.project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    result = _parse_document_internal(document)
    
    if 'error' in result:
        return jsonify(result), 500
    
    return jsonify(result), 200


def _parse_document_internal(document):
    """Internal function to parse document and extract questions."""
    import tempfile
    from datetime import datetime
    from ..services.document_service import DocumentService
    from ..services.extraction_service import QuestionExtractor
    from ..models import Question
    
    # Update status
    document.status = 'processing'
    db.session.commit()
    
    try:
        # Write file data to temp file for processing
        temp_file_path = None
        temp_file_created = False
        
        if document.file_data:
            # Legacy: file stored directly in database
            temp_file = tempfile.NamedTemporaryFile(
                delete=False,
                suffix=f'.{document.file_type}'
            )
            temp_file.write(document.file_data)
            temp_file.close()
            temp_file_path = temp_file.name
            temp_file_created = True
        elif (document.storage_type == 'gcp' or 
              (document.file_metadata and document.file_metadata.get('storage', {}).get('storage_type') == 'gcp')) and document.file_id:
            # GCP Storage: download from cloud storage
            try:
                from app.services.storage_service import get_storage_service
                storage = get_storage_service()
                current_app.logger.info(f"Attempting to download document {document.id} from GCP, file_id: {document.file_id}")
                file_content, metadata = storage.download(document.file_id)
                
                temp_file = tempfile.NamedTemporaryFile(
                    delete=False,
                    suffix=f'.{document.file_type}'
                )
                temp_file.write(file_content)
                temp_file.close()
                temp_file_path = temp_file.name
                temp_file_created = True
                current_app.logger.info(f"Downloaded document {document.id} from GCP storage, size: {len(file_content)} bytes")
            except Exception as e:
                current_app.logger.error(f"Failed to download from GCP: {e}")
                document.status = 'failed'
                document.error_message = f'Failed to download from cloud storage: {str(e)}'
                db.session.commit()
                return {'error': f'Cloud storage download failed: {str(e)}'}
        elif document.storage_type == 'local' and document.file_id:
            # Local storage service: get the local path
            try:
                from app.services.storage_service import get_storage_service
                storage = get_storage_service()
                temp_file_path = storage.get_local_path(document.file_id)
                current_app.logger.info(f"Using local storage path for document {document.id}: {temp_file_path}")
            except Exception as e:
                current_app.logger.error(f"Failed to get local storage path: {e}")
        elif document.file_path:
            # Legacy: direct file path
            temp_file_path = document.file_path
            current_app.logger.info(f"Using legacy file_path for document {document.id}: {temp_file_path}")
        elif document.file_metadata and document.file_metadata.get('storage', {}).get('local_path'):
            # Legacy: storage service path from metadata
            storage_path = document.file_metadata['storage']['local_path']
            import os as path_os
            if path_os.path.exists(storage_path):
                temp_file_path = storage_path
                current_app.logger.info(f"Using file_metadata storage path for document {document.id}: {temp_file_path}")
        
        if not temp_file_path:
            current_app.logger.error(f"No file data available for document {document.id}. storage_type={document.storage_type}, file_id={document.file_id}, has_file_data={document.file_data is not None}")
            document.status = 'failed'
            document.error_message = 'No file data available'
            db.session.commit()
            return {'error': 'No file data available'}

        
        # Extract text from document
        doc_service = DocumentService()
        extracted_text = doc_service.extract_text(temp_file_path, document.file_type)
        
        # Clean up temp file if we created one
        if temp_file_created and temp_file_path:
            import os as temp_os
            try:
                temp_os.unlink(temp_file_path)
            except:
                pass
        
        if not extracted_text:
            document.status = 'failed'
            document.error_message = 'Could not extract text from document'
            db.session.commit()
            return {'error': 'Text extraction failed'}
        
        # Save extracted text
        document.extracted_text = extracted_text
        document.file_metadata = {
            'word_count': len(extracted_text.split()),
            'char_count': len(extracted_text),
        }
        
        # Extract questions
        extractor = QuestionExtractor()
        questions_data = extractor.extract_questions(extracted_text, use_ai=True)
        
        # Question category classification keywords
        category_keywords = {
            'security': ['security', 'encryption', 'authentication', 'password', 'access control', 
                        'firewall', 'vulnerability', 'penetration', 'gdpr', 'hipaa', 'soc', 'iso 27001',
                        'data protection', 'privacy', 'compliance', 'audit'],
            'technical': ['technical', 'architecture', 'infrastructure', 'api', 'integration',
                         'database', 'platform', 'cloud', 'mobile', 'system', 'software', 'hardware'],
            'pricing': ['pricing', 'cost', 'fee', 'budget', 'payment', 'license', 'commercial'],
            'implementation': ['implementation', 'timeline', 'schedule', 'phase', 'milestone', 
                              'deployment', 'go-live', 'rollout', 'project plan'],
            'support': ['support', 'maintenance', 'sla', 'service level', 'helpdesk', 'availability'],
            'team': ['team', 'staff', 'resource', 'personnel', 'experience', 'qualification', 'resume'],
            'training': ['training', 'documentation', 'manual', 'knowledge transfer', 'user guide'],
            'references': ['reference', 'case study', 'past performance', 'similar project', 'client'],
        }
        
        # Create Question records with category classification
        for q_data in questions_data:
            # Classify question category based on content
            q_text_lower = q_data['text'].lower()
            detected_category = 'general'
            detected_section = q_data.get('section', 'Q&A / Questionnaire')
            
            for category, keywords in category_keywords.items():
                if any(kw in q_text_lower for kw in keywords):
                    detected_category = category
                    # Map category to appropriate section name  
                    section_map = {
                        'security': 'Security & Compliance',
                        'technical': 'Technical Approach',
                        'pricing': 'Pricing & Commercial',
                        'implementation': 'Implementation Plan',
                        'support': 'Support & Maintenance',
                        'team': 'Team & Qualifications',
                        'training': 'Training & Documentation',
                        'references': 'References & Experience',
                    }
                    detected_section = section_map.get(category, detected_section)
                    break
            
            question = Question(
                text=q_data['text'],
                section=detected_section,
                category=detected_category,
                order=q_data.get('order', 0),
                status='pending',
                project_id=document.project_id,
                document_id=document.id
            )
            db.session.add(question)

        
        # Update document status
        document.status = 'completed'
        document.processed_at = datetime.utcnow()
        db.session.commit()
        
        # AUTO-ANALYZE RFP AND CREATE SECTIONS
        analysis_result = None
        sections_created = []
        try:
            from ..services.rfp_analysis_agent import RFPAnalysisAgent
            from ..models import RFPSectionType
            
            agent = RFPAnalysisAgent()
            
            # 1. Analyze the RFP document
            analysis_result = agent.analyze_rfp(document.id)
            
            if analysis_result and not analysis_result.get('error'):
                # 2. Get recommended section types
                recommended_sections = analysis_result.get('recommended_sections', [])
                
                if recommended_sections:
                    # Map slugs to section type IDs
                    section_type_ids = []
                    for slug in recommended_sections:
                        section_type = RFPSectionType.query.filter_by(slug=slug, is_active=True).first()
                        if section_type:
                            section_type_ids.append(section_type.id)
                    
                    if section_type_ids:
                        # 3. Auto-create sections with AI content generation
                        created_sections = agent.auto_create_sections(
                            project_id=document.project_id,
                            section_type_ids=section_type_ids,
                            with_generation=True,
                            document_id=document.id
                        )
                        sections_created = created_sections
        except Exception as analysis_error:
            # Log but don't fail - document is still processed
            print(f"Auto-analysis error: {analysis_error}")
            analysis_result = {'error': str(analysis_error)}
        
        # Update project status and completion after successful analysis
        project = document.project
        if project:
            # Update status from draft to in_progress
            if project.status == 'draft':
                project.status = 'in_progress'
            
            # Calculate completion based on questions answered and sections created
            completion = project.calculate_completion()
            project.completion_percent = completion
            
            db.session.commit()
        
        return {
            'message': 'Document parsed and analyzed successfully',
            'document': document.to_dict(),
            'questions_extracted': len(questions_data),
            'analysis': analysis_result,
            'sections_created': len(sections_created) if sections_created else 0
        }
        
    except Exception as e:
        document.status = 'failed'
        document.error_message = str(e)
        db.session.commit()
        return {'error': f'Processing failed: {str(e)}'}


@bp.route('/<int:document_id>', methods=['DELETE'])
@jwt_required()
def delete_document(document_id):
    """Delete a document."""
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    if user.role not in ['admin', 'editor']:
        return jsonify({'error': 'Insufficient permissions'}), 403
    
    document = Document.query.get(document_id)
    
    if not document:
        return jsonify({'error': 'Document not found'}), 404
    
    if document.project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    # Delete file
    if document.file_path and os.path.exists(document.file_path):
        os.remove(document.file_path)
    
    # Delete associated chat sessions first (foreign key constraint)
    try:
        from ..models.document_chat import DocumentChatSession
        DocumentChatSession.query.filter_by(document_id=document_id).delete()
    except Exception as e:
        current_app.logger.warning(f"Could not delete chat sessions: {e}")
    
    # Delete from cloud storage if applicable
    try:
        if document.storage_type == 'gcp' and document.file_id:
            from app.services.storage_service import get_storage_service
            storage = get_storage_service()
            storage.delete(document.file_id)
    except Exception as e:
        current_app.logger.warning(f"Could not delete from cloud storage: {e}")
    
    db.session.delete(document)
    db.session.commit()
    
    return jsonify({'message': 'Document deleted'}), 200


@bp.route('/<int:document_id>/analyze', methods=['POST'])
@jwt_required()
def analyze_document(document_id):
    """Analyze an RFP document to extract structure and suggest proposal sections."""
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    document = Document.query.get(document_id)
    
    if not document:
        return jsonify({'error': 'Document not found'}), 404
    
    if document.project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    # Check if document has been parsed
    if document.status != 'completed' or not document.extracted_text:
        # Auto-parse if needed
        result = _parse_document_internal(document)
        if 'error' in result:
            return jsonify(result), 500
    
    # Run analysis
    from ..services.rfp_analysis_agent import get_rfp_analysis_agent
    agent = get_rfp_analysis_agent()
    analysis = agent.analyze_rfp(document_id)
    
    if 'error' in analysis:
        return jsonify(analysis), 500
    
    return jsonify({
        'message': 'RFP analysis complete',
        **analysis
    }), 200


@bp.route('/<int:document_id>/auto-build', methods=['POST'])
@jwt_required()
def auto_build_proposal(document_id):
    """Automatically create proposal sections based on RFP analysis."""
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    document = Document.query.get(document_id)
    
    if not document:
        return jsonify({'error': 'Document not found'}), 404
    
    if document.project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    data = request.get_json() or {}
    section_type_ids = data.get('section_type_ids', [])
    
    if not section_type_ids:
        return jsonify({'error': 'No section types provided'}), 400
    
    # Create sections with content generation
    from ..services.rfp_analysis_agent import get_rfp_analysis_agent
    agent = get_rfp_analysis_agent()
    result = agent.auto_create_sections(
        project_id=document.project_id,
        section_type_ids=section_type_ids,
        with_generation=data.get('generate_content', True),  # Default to True
        document_id=document_id  # Pass document for context
    )
    
    if 'error' in result:
        return jsonify(result), 500
    
    return jsonify({
        'message': 'Proposal sections created with content',
        **result
    }), 201


@bp.route('/<int:document_id>/signed-url', methods=['GET'])
@jwt_required()
def get_document_signed_url(document_id):
    """
    Get a signed public URL for Microsoft Office Online Viewer.
    
    For GCP-stored documents, generates a signed URL that can be used with:
    - Microsoft Office Viewer: https://view.officeapps.live.com/op/embed.aspx?src={url}
    - Google Docs Viewer: https://docs.google.com/gview?url={url}&embedded=true
    """
    from datetime import timedelta
    from urllib.parse import quote
    
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    document = Document.query.get(document_id)
    
    if not document:
        return jsonify({'error': 'Document not found'}), 404
    
    if document.project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    # Check if document is stored in GCP
    if document.storage_type != 'gcp' or not document.file_id:
        return jsonify({
            'error': 'Document is not stored in cloud storage',
            'message': 'Microsoft Office viewer requires cloud-stored documents'
        }), 400
    
    try:
        from google.cloud import storage as gcs
        
        bucket_name = os.environ.get('GCP_STORAGE_BUCKET') or os.environ.get('GOOGLE_CLOUD_BUCKET_NAME')
        creds_path = os.environ.get('GCP_STORAGE_CREDENTIALS') or os.environ.get('GOOGLE_APPLICATION_CREDENTIALS')
        
        if not bucket_name:
            return jsonify({'error': 'GCP bucket not configured'}), 500
        
        if creds_path and os.path.exists(creds_path):
            client = gcs.Client.from_service_account_json(creds_path)
        else:
            client = gcs.Client()
        
        bucket = client.bucket(bucket_name)
        
        # Get blob using file metadata
        blob_name = None
        if document.file_metadata and document.file_metadata.get('storage', {}).get('blob_name'):
            blob_name = document.file_metadata['storage']['blob_name']
        
        if not blob_name:
            # Try to find blob by file_id
            from app.services.storage_service import get_storage_service
            storage = get_storage_service()
            blob = storage.provider._find_blob(document.file_id)
            if blob:
                blob_name = blob.name
        
        if not blob_name:
            return jsonify({'error': 'Unable to locate file in GCP storage'}), 404
        
        blob = bucket.blob(blob_name)
        
        # Verify blob exists
        if not blob.exists():
            return jsonify({'error': 'File not found in GCP storage'}), 404
        
        # Generate signed URL (valid for 1 hour)
        signed_url = blob.generate_signed_url(
            version="v4",
            expiration=timedelta(hours=1),
            method="GET"
        )
        
        # Build viewer URLs
        microsoft_viewer_url = f"https://view.officeapps.live.com/op/embed.aspx?src={quote(signed_url, safe='')}"
        google_viewer_url = f"https://docs.google.com/gview?url={quote(signed_url, safe='')}&embedded=true"
        
        return jsonify({
            'signed_url': signed_url,
            'microsoft_viewer_url': microsoft_viewer_url,
            'google_viewer_url': google_viewer_url,
            'file_name': document.original_filename,
            'file_type': document.file_type,
            'expires_in': 3600
        }), 200
        
    except ImportError:
        return jsonify({'error': 'GCS library not installed'}), 500
    except Exception as e:
        current_app.logger.error(f"Failed to generate signed URL for document {document_id}: {e}")
        return jsonify({'error': f'Failed to generate signed URL: {str(e)}'}), 500

@bp.route('/<int:document_id>/preview', methods=['GET'])
@jwt_required(optional=True, locations=['headers', 'query_string'])
def preview_document(document_id):
    """Get document content for preview/viewing.
    
    Supports both header-based and query string JWT for iframe embedding.
    """
    from flask import Response, send_file
    import tempfile
    import io
    
    # Get user from JWT (supports both header and query_string)
    user_id = get_jwt_identity()
    if not user_id:
        return jsonify({'error': 'Authentication required'}), 401
    
    user_id = int(user_id)
    user = User.query.get(user_id)
    
    document = Document.query.get(document_id)
    
    if not document:
        return jsonify({'error': 'Document not found'}), 404
    
    if document.project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    # Get file data from database or storage service
    file_data = document.file_data
    
    if not file_data:
        # Try to read from storage service (local path or cloud)
        try:
            from app.services.storage_service import get_storage_service
            storage_service = get_storage_service()
            
            current_app.logger.info(f"Preview: document {document_id}, storage_type={document.storage_type}, file_id={document.file_id}")
            
            if storage_service and document.file_id:
                # Use storage service download method
                current_app.logger.info(f"Attempting to download from storage service...")
                file_bytes, metadata = storage_service.download(document.file_id)
                file_data = file_bytes
                current_app.logger.info(f"Downloaded {len(file_data)} bytes from storage")
            else:
                # Fallback: check local path in metadata
                storage_info = document.file_metadata.get('storage', {}) if document.file_metadata else {}
                local_path = storage_info.get('local_path')
                
                if local_path and os.path.exists(local_path):
                    with open(local_path, 'rb') as f:
                        file_data = f.read()
                    current_app.logger.info(f"Read from local path: {local_path}")
        except Exception as e:
            current_app.logger.error(f"Failed to read file from storage: {e}", exc_info=True)
            return jsonify({'error': f'Failed to load file: {str(e)}'}), 500
    
    if not file_data:
        current_app.logger.error(f"No file data available for document {document_id}")
        return jsonify({'error': 'No file data available'}), 404
    
    file_type = document.file_type.lower()
    
    # PDF - return directly
    if file_type == 'pdf':
        return Response(
            file_data,
            mimetype='application/pdf',
            headers={
                'Content-Disposition': f'inline; filename="{document.original_filename}"'
            }
        )
    
    # DOCX - convert to HTML
    if file_type in ['docx', 'doc']:
        try:
            import mammoth
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_type}')
            temp_file.write(file_data)
            temp_file.close()
            
            with open(temp_file.name, 'rb') as f:
                result = mammoth.convert_to_html(f)
                html_content = result.value
            
            os.unlink(temp_file.name)
            
            # Wrap in styled HTML
            styled_html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; padding: 20px; max-width: 800px; margin: 0 auto; line-height: 1.6; }}
        h1, h2, h3 {{ color: #1f2937; }}
        p {{ margin: 0.5em 0; }}
        table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
        th, td {{ border: 1px solid #e5e7eb; padding: 8px; text-align: left; }}
        th {{ background: #f3f4f6; }}
    </style>
</head>
<body>{html_content}</body>
</html>'''
            
            return Response(styled_html, mimetype='text/html')
            
        except Exception as e:
            return jsonify({'error': f'Failed to convert DOCX: {str(e)}'}), 500
    
    # XLSX - convert to HTML table
    if file_type in ['xlsx', 'xls']:
        try:
            from openpyxl import load_workbook
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_type}')
            temp_file.write(file_data)
            temp_file.close()
            
            wb = load_workbook(temp_file.name)
            os.unlink(temp_file.name)
            
            html_tables = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                html = f'<h2>{sheet_name}</h2><table>'
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True), 1):
                    if row_idx == 1:
                        html += '<thead><tr>'
                        html += ''.join(f'<th>{cell if cell else ""}</th>' for cell in row)
                        html += '</tr></thead><tbody>'
                    else:
                        html += '<tr>'
                        html += ''.join(f'<td>{cell if cell else ""}</td>' for cell in row)
                        html += '</tr>'
                html += '</tbody></table>'
                html_tables.append(html)
            
            styled_html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; padding: 20px; }}
        h2 {{ color: #1f2937; margin-top: 2em; }}
        table {{ border-collapse: collapse; width: 100%; margin: 1em 0; font-size: 14px; }}
        th, td {{ border: 1px solid #e5e7eb; padding: 8px; text-align: left; }}
        th {{ background: #f3f4f6; font-weight: 600; }}
        tr:hover {{ background: #f9fafb; }}
    </style>
</head>
<body>{''.join(html_tables)}</body>
</html>'''
            
            return Response(styled_html, mimetype='text/html')
            
        except Exception as e:
            return jsonify({'error': f'Failed to convert XLSX: {str(e)}'}), 500
    
    # PPTX - convert to HTML with slide content
    if file_type in ['pptx', 'ppt']:
        try:
            from pptx import Presentation
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_type}')
            temp_file.write(file_data)
            temp_file.close()
            
            prs = Presentation(temp_file.name)
            os.unlink(temp_file.name)
            
            slides_html = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(f'<p>{shape.text}</p>')
                
                slides_html.append(f'''
                <div class="slide">
                    <div class="slide-header">Slide {slide_idx}</div>
                    <div class="slide-content">{''.join(slide_content) if slide_content else '<p class="empty">No text content</p>'}</div>
                </div>
                ''')
            
            styled_html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; padding: 20px; background: #f3f4f6; }}
        .slide {{ background: white; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); margin-bottom: 20px; overflow: hidden; }}
        .slide-header {{ background: #3b82f6; color: white; padding: 10px 16px; font-weight: 600; }}
        .slide-content {{ padding: 20px; min-height: 150px; }}
        .slide-content p {{ margin: 0.5em 0; }}
        .empty {{ color: #9ca3af; font-style: italic; }}
    </style>
</head>
<body>{''.join(slides_html)}</body>
</html>'''
            
            return Response(styled_html, mimetype='text/html')
            
        except Exception as e:
            return jsonify({'error': f'Failed to convert PPTX: {str(e)}'}), 500
    
    # Unsupported type
    return jsonify({'error': f'Preview not supported for {file_type}'}), 400


@bp.route('/<int:document_id>/download', methods=['GET'])
@jwt_required()
def download_document(document_id):
    """Download the original document file."""
    from flask import Response
    
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    document = Document.query.get(document_id)
    
    if not document:
        return jsonify({'error': 'Document not found'}), 404
    
    if document.project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    if not document.file_data:
        return jsonify({'error': 'No file data available'}), 404
    
    # Determine MIME type
    mime_types = {
        'pdf': 'application/pdf',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'doc': 'application/msword',
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'xls': 'application/vnd.ms-excel',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'ppt': 'application/vnd.ms-powerpoint',
    }
    
    mime_type = mime_types.get(document.file_type.lower(), 'application/octet-stream')
    
    return Response(
        document.file_data,
        mimetype=mime_type,
        headers={
            'Content-Disposition': f'attachment; filename="{document.original_filename}"'
        }
    )


@bp.route('/search', methods=['POST'])
@jwt_required()
def hybrid_search_documents():
    """
    Hybrid search across document chunks.
    
    Combines dense (semantic) and sparse (BM25 keyword) vectors
    for comprehensive search results using Reciprocal Rank Fusion.
    
    Request body:
        {
            "query": "search query text",
            "limit": 10,  // optional, default 10
            "file_id": "uuid",  // optional, filter by specific document
            "project_id": 123  // optional, filter by project (NOT YET IMPLEMENTED)
        }
    
    Returns:
        {
            "results": [
                {
                    "chunk_id": "abc123",
                    "file_id": "uuid",
                    "page_number": 5,
                    "content": "matched text...",
                    "score": 0.95,
                    "doc_url": "gs://bucket/path",
                    "original_filename": "document.pdf",
                    "metadata": {...}
                }
            ],
            "query": "original query",
            "total_results": 10
        }
    """
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    data = request.get_json()
    if not data:
        return jsonify({'error': 'Request body required'}), 400
    
    query = data.get('query', '').strip()
    if not query:
        return jsonify({'error': 'Query is required'}), 400
    
    limit = data.get('limit', 10)
    file_id = data.get('file_id')
    
    try:
        from app.services.hybrid_search_service import get_hybrid_search_service
        
        hybrid_search = get_hybrid_search_service(user.organization_id)
        
        if not hybrid_search.enabled:
            return jsonify({
                'error': 'Vector search not available',
                'results': [],
                'query': query
            }), 503
        
        results = hybrid_search.hybrid_search(
            query=query,
            org_id=user.organization_id,
            limit=limit,
            file_id=file_id
        )
        
        return jsonify({
            'results': [r.to_dict() for r in results],
            'query': query,
            'total_results': len(results)
        }), 200
        
    except Exception as e:
        current_app.logger.error(f"Hybrid search failed: {e}")
        return jsonify({
            'error': f'Search failed: {str(e)}',
            'results': [],
            'query': query
        }), 500


@bp.route('/<int:document_id>/chunks', methods=['GET'])
@jwt_required()
def get_document_chunks(document_id):
    """
    Get all indexed chunks for a document.
    
    Returns:
        {
            "chunks": [...],
            "total": 25,
            "document_id": 123
        }
    """
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    document = Document.query.get(document_id)
    if not document:
        return jsonify({'error': 'Document not found'}), 404
    
    project = Project.query.get(document.project_id)
    if project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    try:
        from app.services.hybrid_search_service import get_hybrid_search_service
        
        hybrid_search = get_hybrid_search_service(user.organization_id)
        
        if not hybrid_search.enabled or not document.file_id:
            return jsonify({
                'chunks': [],
                'total': 0,
                'document_id': document_id
            })
        
        chunks = hybrid_search.get_document_chunks(
            file_id=document.file_id,
            org_id=user.organization_id
        )
        
        return jsonify({
            'chunks': [c.to_dict() for c in chunks],
            'total': len(chunks),
            'document_id': document_id
        }), 200
        
    except Exception as e:
        current_app.logger.error(f"Failed to get chunks for document {document_id}: {e}")
        return jsonify({'error': str(e)}), 500


@bp.route('/<int:document_id>/reindex', methods=['POST'])
@jwt_required()
def reindex_document(document_id):
    """
    Trigger re-indexing of a document's embeddings.
    
    Useful when document chunking or embedding has failed
    and needs to be retried.
    """
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    document = Document.query.get(document_id)
    if not document:
        return jsonify({'error': 'Document not found'}), 404
    
    project = Project.query.get(document.project_id)
    if project.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    try:
        from app.tasks import reprocess_document_embeddings_task
        
        reprocess_document_embeddings_task.delay(document.id, user.organization_id)
        
        document.embedding_status = 'pending'
        db.session.commit()
        
        return jsonify({
            'message': 'Document reindexing triggered',
            'document': document.to_dict()
        }), 200
        
    except Exception as e:
        current_app.logger.error(f"Failed to trigger reindex for document {document_id}: {e}")
        return jsonify({'error': str(e)}), 500
