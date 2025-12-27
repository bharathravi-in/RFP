"""
File preview endpoint for knowledge items.
Serves file content from database (file_data column).
"""
import os
import tempfile
from io import BytesIO
from flask import Blueprint, request, jsonify, send_file, current_app, Response
from flask_jwt_extended import jwt_required, get_jwt_identity
from ..models import KnowledgeItem, User

bp = Blueprint('preview', __name__)


@bp.route('/<int:item_id>', methods=['GET'])
@jwt_required()
def preview_file(item_id):
    """Get file content for preview."""
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    item = KnowledgeItem.query.get(item_id)
    
    if not item:
        return jsonify({'error': 'Item not found'}), 404
    
    if item.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    # For manual entries, return content directly
    if item.source_type == 'manual' or not item.file_data:
        return jsonify({
            'type': 'text',
            'title': item.title,
            'content': item.content,
            'metadata': item.item_metadata
        }), 200
    
    file_type = item.file_type or ''
    
    # For text files, decode and return content
    if file_type.startswith('text/') or (item.source_file and item.source_file.endswith(('.txt', '.md', '.csv'))):
        try:
            content = item.file_data.decode('utf-8')
            return jsonify({
                'type': 'text',
                'title': item.title,
                'content': content,
                'file_type': file_type,
                'can_download': True
            }), 200
        except Exception as e:
            # Fall back to extracted content
            return jsonify({
                'type': 'document',
                'title': item.title,
                'content': item.content,
                'file_type': file_type,
                'file_name': item.source_file,
                'can_download': True
            }), 200
    
    # For PDFs and other documents, return metadata + extracted text
    return jsonify({
        'type': 'document',
        'title': item.title,
        'content': item.content,  # Extracted text
        'file_type': file_type,
        'file_name': item.source_file,
        'file_size': item.file_size,
        'can_download': True
    }), 200


@bp.route('/<int:item_id>/download', methods=['GET'])
@jwt_required()
def download_file(item_id):
    """Download the original file from database."""
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    item = KnowledgeItem.query.get(item_id)
    
    if not item:
        return jsonify({'error': 'Item not found'}), 404
    
    if item.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    if not item.file_data:
        return jsonify({'error': 'File not available'}), 404
    
    # Create BytesIO from file_data
    file_buffer = BytesIO(item.file_data)
    
    return send_file(
        file_buffer,
        mimetype=item.file_type or 'application/octet-stream',
        as_attachment=True,
        download_name=item.source_file or item.title
    )


@bp.route('/<int:item_id>/signed-url', methods=['GET'])
@jwt_required()
def get_signed_url(item_id):
    """
    Get a signed public URL for Microsoft Office Online Viewer.
    
    This endpoint uploads the file to Google Cloud Storage (if not already there)
    and returns a signed URL that can be used with Microsoft Office Online Viewer
    or Google Docs Viewer for proper preview with all features (e.g., Excel sheet tabs).
    
    Returns:
        JSON with 'url' for Microsoft viewer: 
        https://view.officeapps.live.com/op/embed.aspx?src={signed_url}
    """
    import os
    from io import BytesIO
    
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    item = KnowledgeItem.query.get(item_id)
    
    if not item:
        return jsonify({'error': 'Item not found'}), 404
    
    if item.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    if not item.file_data:
        return jsonify({'error': 'File not available'}), 404
    
    # Check if GCS bucket is configured (works even if main storage is local)
    bucket_name = os.environ.get('GCP_STORAGE_BUCKET') or os.environ.get('GOOGLE_CLOUD_BUCKET_NAME')
    
    current_app.logger.info(f"Signed URL request for item {item_id}, bucket: {bucket_name}")
    
    if not bucket_name:
        return jsonify({
            'error': 'Cloud storage not configured',
            'message': 'Microsoft Office Viewer requires GCS bucket. Set GOOGLE_CLOUD_BUCKET_NAME or GCP_STORAGE_BUCKET environment variable.'
        }), 400
    
    try:
        from google.cloud import storage as gcs
        
        # Get credentials
        creds_path = os.environ.get('GCP_STORAGE_CREDENTIALS') or os.environ.get('GOOGLE_APPLICATION_CREDENTIALS')
        project_id = os.environ.get('GOOGLE_CLOUD_PROJECT_ID')
        
        if creds_path and os.path.exists(creds_path):
            client = gcs.Client.from_service_account_json(creds_path)
        else:
            client = gcs.Client(project=project_id)
        
        bucket = client.bucket(bucket_name)
        
        # Create blob path for preview files
        from datetime import datetime
        date_prefix = datetime.utcnow().strftime('%Y/%m/%d')
        file_ext = item.source_file.rsplit('.', 1)[-1] if item.source_file and '.' in item.source_file else 'bin'
        blob_name = f"previews/{date_prefix}/{item_id}_{item.source_file or 'file'}"
        
        blob = bucket.blob(blob_name)
        
        # Upload file data to GCS
        content_type = item.file_type or 'application/octet-stream'
        blob.upload_from_string(item.file_data, content_type=content_type)
        
        # Generate signed URL (valid for 1 hour)
        from datetime import timedelta
        signed_url = blob.generate_signed_url(
            version="v4",
            expiration=timedelta(hours=1),
            method="GET"
        )
        
        # Build Microsoft viewer URL
        from urllib.parse import quote
        microsoft_viewer_url = f"https://view.officeapps.live.com/op/embed.aspx?src={quote(signed_url, safe='')}"
        google_viewer_url = f"https://docs.google.com/gview?url={quote(signed_url, safe='')}&embedded=true"
        
        return jsonify({
            'signed_url': signed_url,
            'microsoft_viewer_url': microsoft_viewer_url,
            'google_viewer_url': google_viewer_url,
            'file_name': item.source_file or item.title,
            'file_type': item.file_type,
            'expires_in': 3600  # 1 hour in seconds
        }), 200
        
    except ImportError:
        return jsonify({
            'error': 'GCS library not installed',
            'message': 'Install google-cloud-storage: pip install google-cloud-storage'
        }), 500
    except Exception as e:
        current_app.logger.error(f"Failed to generate signed URL: {e}")
        return jsonify({
            'error': 'Failed to generate signed URL',
            'message': str(e)
        }), 500


@bp.route('/<int:item_id>/file', methods=['GET'])
@jwt_required()
def serve_file(item_id):
    """Serve file for inline viewing (PDF preview in iframe)."""
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    item = KnowledgeItem.query.get(item_id)
    
    if not item:
        return jsonify({'error': 'Item not found'}), 404
    
    if item.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    if not item.file_data:
        return jsonify({'error': 'File not available'}), 404
    
    # Determine MIME type
    mime_type = item.file_type or 'application/octet-stream'
    if item.source_file and item.source_file.endswith('.pdf'):
        mime_type = 'application/pdf'
    
    # Create BytesIO from file_data
    file_buffer = BytesIO(item.file_data)
    
    return send_file(
        file_buffer,
        mimetype=mime_type,
        as_attachment=False  # Inline display
    )


@bp.route('/<int:item_id>/view', methods=['GET'])
@jwt_required()
def view_file(item_id):
    """View file with proper HTML rendering for different file types.
    
    Similar to RFP document preview - converts DOCX, XLSX, PPTX to HTML.
    """
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    item = KnowledgeItem.query.get(item_id)
    
    if not item:
        return jsonify({'error': 'Item not found'}), 404
    
    if item.organization_id != user.organization_id:
        return jsonify({'error': 'Access denied'}), 403
    
    if not item.file_data:
        return jsonify({'error': 'No file data available'}), 404
    
    # Get file extension from source_file or file_type
    file_ext = ''
    if item.source_file:
        file_ext = item.source_file.rsplit('.', 1)[-1].lower() if '.' in item.source_file else ''
    elif item.file_type:
        # Handle mime types like 'application/pdf'
        if '/' in item.file_type:
            file_ext = item.file_type.split('/')[-1].lower()
        else:
            file_ext = item.file_type.lower()
    
    # PDF - return directly
    if file_ext == 'pdf':
        return Response(
            item.file_data,
            mimetype='application/pdf',
            headers={
                'Content-Disposition': f'inline; filename="{item.source_file or item.title}"'
            }
        )
    
    # DOCX - convert to HTML
    if file_ext in ['docx', 'doc']:
        try:
            import mammoth
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_ext}')
            temp_file.write(item.file_data)
            temp_file.close()
            
            with open(temp_file.name, 'rb') as f:
                result = mammoth.convert_to_html(f)
                html_content = result.value
            
            os.unlink(temp_file.name)
            
            styled_html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; padding: 20px; max-width: 800px; margin: 0 auto; line-height: 1.6; }}
        h1, h2, h3 {{ color: #1f2937; }}
        p {{ margin: 0.5em 0; }}
        table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
        th, td {{ border: 1px solid #e5e7eb; padding: 8px; text-align: left; }}
        th {{ background: #f3f4f6; }}
    </style>
</head>
<body>{html_content}</body>
</html>'''
            
            return Response(styled_html, mimetype='text/html')
            
        except Exception as e:
            return jsonify({'error': f'Failed to convert DOCX: {str(e)}'}), 500
    
    # XLSX - convert to HTML table
    if file_ext in ['xlsx', 'xls']:
        try:
            from openpyxl import load_workbook
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_ext}')
            temp_file.write(item.file_data)
            temp_file.close()
            
            wb = load_workbook(temp_file.name)
            os.unlink(temp_file.name)
            
            html_tables = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                html = f'<h2>{sheet_name}</h2><table>'
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True), 1):
                    if row_idx == 1:
                        html += '<thead><tr>'
                        html += ''.join(f'<th>{cell if cell else ""}</th>' for cell in row)
                        html += '</tr></thead><tbody>'
                    else:
                        html += '<tr>'
                        html += ''.join(f'<td>{cell if cell else ""}</td>' for cell in row)
                        html += '</tr>'
                html += '</tbody></table>'
                html_tables.append(html)
            
            styled_html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; padding: 20px; }}
        h2 {{ color: #1f2937; margin-top: 2em; }}
        table {{ border-collapse: collapse; width: 100%; margin: 1em 0; font-size: 14px; }}
        th, td {{ border: 1px solid #e5e7eb; padding: 8px; text-align: left; }}
        th {{ background: #f3f4f6; font-weight: 600; }}
        tr:hover {{ background: #f9fafb; }}
    </style>
</head>
<body>{''.join(html_tables)}</body>
</html>'''
            
            return Response(styled_html, mimetype='text/html')
            
        except Exception as e:
            return jsonify({'error': f'Failed to convert XLSX: {str(e)}'}), 500
    
    # PPTX - convert to HTML with slide content
    if file_ext in ['pptx', 'ppt']:
        try:
            from pptx import Presentation
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_ext}')
            temp_file.write(item.file_data)
            temp_file.close()
            
            prs = Presentation(temp_file.name)
            os.unlink(temp_file.name)
            
            slides_html = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(f'<p>{shape.text}</p>')
                
                slides_html.append(f'''
                <div class="slide">
                    <div class="slide-header">Slide {slide_idx}</div>
                    <div class="slide-content">{''.join(slide_content) if slide_content else '<p class="empty">No text content</p>'}</div>
                </div>
                ''')
            
            styled_html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; padding: 20px; background: #f3f4f6; }}
        .slide {{ background: white; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); margin-bottom: 20px; overflow: hidden; }}
        .slide-header {{ background: #3b82f6; color: white; padding: 10px 16px; font-weight: 600; }}
        .slide-content {{ padding: 20px; min-height: 150px; }}
        .slide-content p {{ margin: 0.5em 0; }}
        .empty {{ color: #9ca3af; font-style: italic; }}
    </style>
</head>
<body>{''.join(slides_html)}</body>
</html>'''
            
            return Response(styled_html, mimetype='text/html')
            
        except Exception as e:
            return jsonify({'error': f'Failed to convert PPTX: {str(e)}'}), 500
    
    # Text files - show as preformatted text
    if file_ext in ['txt', 'md', 'csv', 'json']:
        try:
            content = item.file_data.decode('utf-8')
            styled_html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>
        body {{ font-family: monospace; padding: 20px; white-space: pre-wrap; word-wrap: break-word; }}
    </style>
</head>
<body>{content}</body>
</html>'''
            return Response(styled_html, mimetype='text/html')
        except:
            pass
    
    # Unsupported type - return extracted text
    if item.content:
        styled_html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; padding: 20px; line-height: 1.6; }}
        .notice {{ background: #fef3c7; border: 1px solid #f59e0b; padding: 12px; border-radius: 8px; margin-bottom: 20px; }}
        .content {{ white-space: pre-wrap; }}
    </style>
</head>
<body>
    <div class="notice">⚠️ Direct preview not available for this file type. Showing extracted text content below.</div>
    <div class="content">{item.content}</div>
</body>
</html>'''
        return Response(styled_html, mimetype='text/html')
    
    return jsonify({'error': f'Preview not supported for this file type'}), 400
