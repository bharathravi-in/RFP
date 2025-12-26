"""
Document Chunking Service with Docling Integration.

Provides intelligent document chunking with page-level extraction
for PDFs, DOCX, and other document formats. Designed for Qdrant hybrid search.
"""

import os
import logging
import hashlib
from typing import List, Dict, Optional, Any
from dataclasses import dataclass, field
from datetime import datetime

logger = logging.getLogger(__name__)


@dataclass
class DocumentChunk:
    """Represents a chunk of a document with rich metadata."""
    
    # Core identification
    chunk_id: str
    file_id: str
    page_number: int
    chunk_index: int  # Index within the page
    
    # Content
    content: str
    content_type: str  # 'text', 'table', 'heading', 'paragraph', 'slide', 'sheet'
    
    # Metrics
    word_count: int = 0
    char_count: int = 0
    sentence_count: int = 0
    
    # Content features
    has_tables: bool = False
    has_images: bool = False
    has_code: bool = False
    headings: List[str] = field(default_factory=list)
    keywords: List[str] = field(default_factory=list)
    
    # Parent document info (denormalized for search)
    doc_url: str = None
    original_filename: str = None
    
    # Status
    status: str = 'active'  # 'active', 'deleted', 'archived'
    
    # Processing metadata
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    # Embeddings (populated separately)
    dense_embedding: List[float] = None
    sparse_indices: List[int] = None
    sparse_values: List[float] = None
    
    def to_dict(self) -> Dict:
        """Convert to dictionary for storage."""
        return {
            'chunk_id': self.chunk_id,
            'file_id': self.file_id,
            'page_number': self.page_number,
            'chunk_index': self.chunk_index,
            'content': self.content,
            'content_type': self.content_type,
            'word_count': self.word_count,
            'char_count': self.char_count,
            'sentence_count': self.sentence_count,
            'has_tables': self.has_tables,
            'has_images': self.has_images,
            'has_code': self.has_code,
            'headings': self.headings,
            'keywords': self.keywords,
            'doc_url': self.doc_url,
            'original_filename': self.original_filename,
            'status': self.status,
            'metadata': self.metadata
        }
    
    def to_qdrant_metadata(self, org_id: int) -> Dict:
        """Get metadata for Qdrant point storage."""
        return {
            'file_id': self.file_id,
            'doc_url': self.doc_url or '',
            'page_number': self.page_number,
            'chunk_index': self.chunk_index,
            'status': self.status,
            'org_id': org_id,
            'content_type': self.content_type,
            'word_count': self.word_count,
            'char_count': self.char_count,
            'has_tables': self.has_tables,
            'has_images': self.has_images,
            'original_filename': self.original_filename or '',
            'headings': self.headings,
            'keywords': self.keywords
        }


@dataclass
class ChunkingResult:
    """Result of document chunking operation."""
    file_id: str
    original_filename: str
    file_type: str
    total_pages: int
    total_chunks: int
    total_words: int
    total_chars: int
    chunks: List[DocumentChunk]
    document_metadata: Dict[str, Any]
    processing_time_ms: float
    doc_url: str = None
    
    def to_dict(self) -> Dict:
        """Convert to dictionary."""
        return {
            'file_id': self.file_id,
            'original_filename': self.original_filename,
            'file_type': self.file_type,
            'total_pages': self.total_pages,
            'total_chunks': self.total_chunks,
            'total_words': self.total_words,
            'total_chars': self.total_chars,
            'chunks': [c.to_dict() for c in self.chunks],
            'document_metadata': self.document_metadata,
            'processing_time_ms': self.processing_time_ms,
            'doc_url': self.doc_url
        }


class DoclingChunkingService:
    """
    Enhanced document chunking service with Docling integration.
    
    Provides:
    - Page-based chunking for PDFs
    - Section-based chunking for DOCX
    - Slide-based chunking for PPTX
    - Sheet-based chunking for XLSX
    - Rich metadata extraction (tables, images, headings)
    - Fallback to basic extraction when Docling unavailable
    """
    
    # Maximum chunk size in characters
    MAX_CHUNK_SIZE = 3000
    CHUNK_OVERLAP = 200
    
    def __init__(self):
        self.docling_available = False
        self._init_docling()
    
    def _init_docling(self):
        """Initialize Docling if available."""
        try:
            from docling.document_converter import DocumentConverter
            self.docling_available = True
            self._converter = None  # Lazy initialization
            logger.info("Docling initialized successfully")
        except ImportError:
            logger.warning("Docling not available, using fallback extraction")
            self.docling_available = False
    
    def _get_converter(self):
        """Get or create Docling converter (lazy initialization)."""
        if self._converter is None:
            from docling.document_converter import DocumentConverter
            self._converter = DocumentConverter()
        return self._converter
    
    def _generate_chunk_id(self, file_id: str, page: int, index: int) -> str:
        """Generate unique chunk ID."""
        raw = f"{file_id}:p{page}:c{index}"
        return hashlib.md5(raw.encode()).hexdigest()[:16]
    
    def _count_sentences(self, text: str) -> int:
        """Count sentences in text."""
        import re
        sentences = re.split(r'[.!?]+', text)
        return len([s for s in sentences if s.strip()])
    
    def _extract_keywords(self, text: str, max_keywords: int = 10) -> List[str]:
        """Extract important keywords from text."""
        import re
        # Simple keyword extraction - words that appear important
        words = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', text)
        # Also get capitalized phrases and technical terms
        tech_terms = re.findall(r'\b[A-Z]{2,}\b', text)
        unique = list(set(words + tech_terms))
        return unique[:max_keywords]
    
    def _detect_code(self, text: str) -> bool:
        """Detect if text contains code snippets."""
        import re
        code_patterns = [
            r'```',  # Markdown code blocks
            r'def\s+\w+\s*\(',  # Python functions
            r'function\s+\w+\s*\(',  # JS functions
            r'class\s+\w+',  # Class definitions
            r'\{\s*\n.*\n\s*\}',  # Code blocks
        ]
        for pattern in code_patterns:
            if re.search(pattern, text):
                return True
        return False
    
    def chunk_document(
        self,
        file_path: str,
        file_id: str,
        doc_url: str = None,
        original_filename: str = None,
        file_type: str = None,
        max_chunk_size: int = None
    ) -> ChunkingResult:
        """
        Chunk a document into page-based segments.
        
        Args:
            file_path: Path to the document file
            file_id: Unique identifier for the document
            doc_url: URL where document is stored
            original_filename: Original name of the file
            file_type: File extension (pdf, docx, etc.)
            max_chunk_size: Max characters per chunk
            
        Returns:
            ChunkingResult with all chunks and metadata
        """
        import time
        start_time = time.time()
        
        if not file_type:
            file_type = file_path.rsplit('.', 1)[-1].lower() if '.' in file_path else ''
        
        if not original_filename:
            original_filename = os.path.basename(file_path)
        
        max_chunk_size = max_chunk_size or self.MAX_CHUNK_SIZE
        
        # Use Docling if available and supported
        if self.docling_available and file_type in ['pdf', 'docx', 'doc', 'pptx', 'ppt']:
            try:
                result = self._chunk_with_docling(
                    file_path, file_id, file_type, doc_url, original_filename
                )
            except Exception as e:
                logger.warning(f"Docling failed, using fallback: {e}")
                result = self._chunk_fallback(
                    file_path, file_id, file_type, doc_url, original_filename, max_chunk_size
                )
        else:
            result = self._chunk_fallback(
                file_path, file_id, file_type, doc_url, original_filename, max_chunk_size
            )
        
        result.processing_time_ms = (time.time() - start_time) * 1000
        logger.info(f"Chunked {file_id}: {result.total_chunks} chunks from {result.total_pages} pages in {result.processing_time_ms:.1f}ms")
        
        return result
    
    def _chunk_with_docling(
        self,
        file_path: str,
        file_id: str,
        file_type: str,
        doc_url: str,
        original_filename: str
    ) -> ChunkingResult:
        """Chunk document using Docling for intelligent extraction."""
        converter = self._get_converter()
        
        # Convert document
        result = converter.convert(file_path)
        doc = result.document
        
        chunks = []
        page_contents = {}  # Group content by page
        total_words = 0
        total_chars = 0
        
        # Extract text items with page info
        for item, level in doc.iterate_items():
            page_num = getattr(item, 'page_no', 1) or 1
            if page_num not in page_contents:
                page_contents[page_num] = {
                    'texts': [],
                    'headings': [],
                    'has_tables': False,
                    'has_images': False
                }
            
            content = getattr(item, 'text', '') or str(item)
            item_type = type(item).__name__.lower()
            
            if 'heading' in item_type or 'title' in item_type:
                page_contents[page_num]['headings'].append(content)
            if 'table' in item_type:
                page_contents[page_num]['has_tables'] = True
            if 'image' in item_type or 'figure' in item_type:
                page_contents[page_num]['has_images'] = True
            
            page_contents[page_num]['texts'].append(content)
        
        # Create chunks per page
        for page_num, page_data in sorted(page_contents.items()):
            page_text = '\n'.join(page_data['texts'])
            
            if not page_text.strip():
                continue
            
            word_count = len(page_text.split())
            char_count = len(page_text)
            total_words += word_count
            total_chars += char_count
            
            # Split large pages into sub-chunks
            if char_count > self.MAX_CHUNK_SIZE:
                page_chunks = self._split_into_chunks(
                    page_text, file_id, page_num, doc_url, original_filename,
                    page_data['headings'], page_data['has_tables'], page_data['has_images']
                )
                chunks.extend(page_chunks)
            else:
                chunk = DocumentChunk(
                    chunk_id=self._generate_chunk_id(file_id, page_num, 0),
                    file_id=file_id,
                    page_number=page_num,
                    chunk_index=0,
                    content=page_text,
                    content_type='page',
                    word_count=word_count,
                    char_count=char_count,
                    sentence_count=self._count_sentences(page_text),
                    has_tables=page_data['has_tables'],
                    has_images=page_data['has_images'],
                    has_code=self._detect_code(page_text),
                    headings=page_data['headings'],
                    keywords=self._extract_keywords(page_text),
                    doc_url=doc_url,
                    original_filename=original_filename,
                    metadata={
                        'extraction_method': 'docling',
                        'file_type': file_type
                    }
                )
                chunks.append(chunk)
        
        # Document metadata
        doc_metadata = {
            'title': getattr(doc, 'title', None),
            'author': getattr(doc, 'author', None),
            'total_pages': len(page_contents),
            'file_type': file_type,
            'extraction_method': 'docling'
        }
        
        return ChunkingResult(
            file_id=file_id,
            original_filename=original_filename,
            file_type=file_type,
            total_pages=len(page_contents),
            total_chunks=len(chunks),
            total_words=total_words,
            total_chars=total_chars,
            chunks=chunks,
            document_metadata=doc_metadata,
            processing_time_ms=0,
            doc_url=doc_url
        )
    
    def _chunk_fallback(
        self,
        file_path: str,
        file_id: str,
        file_type: str,
        doc_url: str,
        original_filename: str,
        max_chunk_size: int
    ) -> ChunkingResult:
        """Fallback chunking using basic extractors."""
        chunks = []
        doc_metadata = {}
        total_words = 0
        total_chars = 0
        total_pages = 0
        
        try:
            if file_type == 'pdf':
                chunks, doc_metadata = self._chunk_pdf(
                    file_path, file_id, doc_url, original_filename, max_chunk_size
                )
            elif file_type in ['docx', 'doc']:
                chunks, doc_metadata = self._chunk_docx(
                    file_path, file_id, doc_url, original_filename, max_chunk_size
                )
            elif file_type in ['pptx', 'ppt']:
                chunks, doc_metadata = self._chunk_pptx(
                    file_path, file_id, doc_url, original_filename
                )
            elif file_type in ['xlsx', 'xls']:
                chunks, doc_metadata = self._chunk_xlsx(
                    file_path, file_id, doc_url, original_filename
                )
            else:
                # Generic text file
                chunks, doc_metadata = self._chunk_text_file(
                    file_path, file_id, doc_url, original_filename, max_chunk_size
                )
        except Exception as e:
            logger.error(f"Error chunking document {file_id}: {e}")
            raise
        
        for chunk in chunks:
            total_words += chunk.word_count
            total_chars += chunk.char_count
        
        total_pages = doc_metadata.get('total_pages', len(chunks))
        
        return ChunkingResult(
            file_id=file_id,
            original_filename=original_filename,
            file_type=file_type,
            total_pages=total_pages,
            total_chunks=len(chunks),
            total_words=total_words,
            total_chars=total_chars,
            chunks=chunks,
            document_metadata=doc_metadata,
            processing_time_ms=0,
            doc_url=doc_url
        )
    
    def _chunk_pdf(
        self,
        file_path: str,
        file_id: str,
        doc_url: str,
        original_filename: str,
        max_chunk_size: int
    ) -> tuple:
        """Chunk PDF by pages using pdfplumber."""
        import pdfplumber
        
        chunks = []
        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ''
                tables = page.extract_tables()
                
                has_tables = bool(tables)
                table_text = ''
                if tables:
                    for table in tables:
                        if table:
                            for row in table:
                                row_text = ' | '.join(str(cell or '') for cell in row)
                                table_text += row_text + '\n'
                
                full_text = text
                if table_text:
                    full_text += '\n\n[TABLES]\n' + table_text
                
                if not full_text.strip():
                    continue
                
                # Split if too large
                if len(full_text) > max_chunk_size:
                    page_chunks = self._split_into_chunks(
                        full_text, file_id, page_num, doc_url, original_filename,
                        [], has_tables, False
                    )
                    chunks.extend(page_chunks)
                else:
                    chunk = DocumentChunk(
                        chunk_id=self._generate_chunk_id(file_id, page_num, 0),
                        file_id=file_id,
                        page_number=page_num,
                        chunk_index=0,
                        content=full_text,
                        content_type='page',
                        word_count=len(full_text.split()),
                        char_count=len(full_text),
                        sentence_count=self._count_sentences(full_text),
                        has_tables=has_tables,
                        keywords=self._extract_keywords(full_text),
                        doc_url=doc_url,
                        original_filename=original_filename,
                        metadata={'extraction_method': 'pdfplumber'}
                    )
                    chunks.append(chunk)
        
        doc_metadata = {
            'total_pages': total_pages,
            'file_type': 'pdf',
            'extraction_method': 'pdfplumber'
        }
        
        return chunks, doc_metadata
    
    def _chunk_docx(
        self,
        file_path: str,
        file_id: str,
        doc_url: str,
        original_filename: str,
        max_chunk_size: int
    ) -> tuple:
        """Chunk DOCX by sections/paragraphs."""
        from docx import Document
        
        doc = Document(file_path)
        chunks = []
        current_page = 1
        current_text = []
        headings = []
        char_count = 0
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            if para.style.name.startswith('Heading'):
                headings.append(text)
            
            current_text.append(text)
            char_count += len(text)
            
            # Create chunk when size limit reached
            if char_count > max_chunk_size:
                content = '\n'.join(current_text)
                chunk = DocumentChunk(
                    chunk_id=self._generate_chunk_id(file_id, current_page, 0),
                    file_id=file_id,
                    page_number=current_page,
                    chunk_index=0,
                    content=content,
                    content_type='section',
                    word_count=len(content.split()),
                    char_count=len(content),
                    sentence_count=self._count_sentences(content),
                    headings=headings.copy(),
                    keywords=self._extract_keywords(content),
                    doc_url=doc_url,
                    original_filename=original_filename,
                    metadata={'extraction_method': 'python-docx'}
                )
                chunks.append(chunk)
                current_text = []
                headings = []
                char_count = 0
                current_page += 1
        
        # Add remaining text
        if current_text:
            content = '\n'.join(current_text)
            chunk = DocumentChunk(
                chunk_id=self._generate_chunk_id(file_id, current_page, 0),
                file_id=file_id,
                page_number=current_page,
                chunk_index=0,
                content=content,
                content_type='section',
                word_count=len(content.split()),
                char_count=len(content),
                sentence_count=self._count_sentences(content),
                headings=headings,
                keywords=self._extract_keywords(content),
                doc_url=doc_url,
                original_filename=original_filename,
                metadata={'extraction_method': 'python-docx'}
            )
            chunks.append(chunk)
        
        doc_metadata = {
            'total_pages': len(chunks),
            'file_type': 'docx',
            'extraction_method': 'python-docx'
        }
        
        return chunks, doc_metadata
    
    def _chunk_pptx(
        self,
        file_path: str,
        file_id: str,
        doc_url: str,
        original_filename: str
    ) -> tuple:
        """Chunk PPTX by slides."""
        from pptx import Presentation
        
        prs = Presentation(file_path)
        chunks = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    texts.append(shape.text)
            
            if texts:
                content = '\n'.join(texts)
                chunk = DocumentChunk(
                    chunk_id=self._generate_chunk_id(file_id, slide_num, 0),
                    file_id=file_id,
                    page_number=slide_num,
                    chunk_index=0,
                    content=content,
                    content_type='slide',
                    word_count=len(content.split()),
                    char_count=len(content),
                    sentence_count=self._count_sentences(content),
                    keywords=self._extract_keywords(content),
                    doc_url=doc_url,
                    original_filename=original_filename,
                    metadata={'extraction_method': 'python-pptx', 'slide_number': slide_num}
                )
                chunks.append(chunk)
        
        doc_metadata = {
            'total_pages': len(prs.slides),
            'file_type': 'pptx',
            'extraction_method': 'python-pptx'
        }
        
        return chunks, doc_metadata
    
    def _chunk_xlsx(
        self,
        file_path: str,
        file_id: str,
        doc_url: str,
        original_filename: str
    ) -> tuple:
        """Chunk XLSX by sheets."""
        from openpyxl import load_workbook
        
        wb = load_workbook(file_path)
        chunks = []
        
        for sheet_num, sheet in enumerate(wb.worksheets, 1):
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) for cell in row if cell)
                if row_text:
                    rows.append(row_text)
            
            if rows:
                content = '\n'.join(rows)
                chunk = DocumentChunk(
                    chunk_id=self._generate_chunk_id(file_id, sheet_num, 0),
                    file_id=file_id,
                    page_number=sheet_num,
                    chunk_index=0,
                    content=content,
                    content_type='sheet',
                    word_count=len(content.split()),
                    char_count=len(content),
                    has_tables=True,
                    doc_url=doc_url,
                    original_filename=original_filename,
                    metadata={
                        'sheet_name': sheet.title,
                        'extraction_method': 'openpyxl'
                    }
                )
                chunks.append(chunk)
        
        doc_metadata = {
            'total_pages': len(wb.worksheets),
            'file_type': 'xlsx',
            'extraction_method': 'openpyxl'
        }
        
        return chunks, doc_metadata
    
    def _chunk_text_file(
        self,
        file_path: str,
        file_id: str,
        doc_url: str,
        original_filename: str,
        max_chunk_size: int
    ) -> tuple:
        """Chunk plain text file."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        
        chunks = self._split_into_chunks(
            text, file_id, 1, doc_url, original_filename, [], False, False
        )
        
        doc_metadata = {
            'total_pages': 1,
            'file_type': 'txt',
            'extraction_method': 'raw'
        }
        
        return chunks, doc_metadata
    
    def _split_into_chunks(
        self,
        text: str,
        file_id: str,
        page_number: int,
        doc_url: str,
        original_filename: str,
        headings: List[str],
        has_tables: bool,
        has_images: bool
    ) -> List[DocumentChunk]:
        """Split large text into overlapping chunks."""
        if len(text) <= self.MAX_CHUNK_SIZE:
            return [DocumentChunk(
                chunk_id=self._generate_chunk_id(file_id, page_number, 0),
                file_id=file_id,
                page_number=page_number,
                chunk_index=0,
                content=text,
                content_type='text',
                word_count=len(text.split()),
                char_count=len(text),
                sentence_count=self._count_sentences(text),
                has_tables=has_tables,
                has_images=has_images,
                headings=headings,
                keywords=self._extract_keywords(text),
                doc_url=doc_url,
                original_filename=original_filename,
                metadata={'is_split': False}
            )]
        
        chunks = []
        start = 0
        chunk_index = 0
        
        while start < len(text):
            end = start + self.MAX_CHUNK_SIZE
            
            # Try to break at sentence boundary
            if end < len(text):
                for sep in ['. ', '.\n', '\n\n', '\n', ' ']:
                    last_sep = text[start:end].rfind(sep)
                    if last_sep > self.MAX_CHUNK_SIZE // 2:
                        end = start + last_sep + len(sep)
                        break
            
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append(DocumentChunk(
                    chunk_id=self._generate_chunk_id(file_id, page_number, chunk_index),
                    file_id=file_id,
                    page_number=page_number,
                    chunk_index=chunk_index,
                    content=chunk_text,
                    content_type='text',
                    word_count=len(chunk_text.split()),
                    char_count=len(chunk_text),
                    sentence_count=self._count_sentences(chunk_text),
                    has_tables=has_tables and chunk_index == 0,
                    has_images=has_images and chunk_index == 0,
                    headings=headings if chunk_index == 0 else [],
                    keywords=self._extract_keywords(chunk_text),
                    doc_url=doc_url,
                    original_filename=original_filename,
                    metadata={'is_split': True, 'split_index': chunk_index}
                ))
                chunk_index += 1
            
            start = end - self.CHUNK_OVERLAP if end < len(text) else len(text)
        
        return chunks


# Singleton instance
_docling_chunking_service: Optional[DoclingChunkingService] = None


def get_docling_chunking_service() -> DoclingChunkingService:
    """Get the Docling chunking service instance."""
    global _docling_chunking_service
    if _docling_chunking_service is None:
        _docling_chunking_service = DoclingChunkingService()
    return _docling_chunking_service
