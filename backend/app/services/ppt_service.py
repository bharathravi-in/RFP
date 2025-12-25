"""
PPT Service - PowerPoint file generation using python-pptx
Creates professional .pptx files from slide content
"""
import io
from typing import Dict, List, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging

logger = logging.getLogger(__name__)


class PPTService:
    """Service for generating PowerPoint presentations."""
    
    # Default branding
    DEFAULT_COLORS = {
        'primary': RGBColor(75, 0, 130),      # Indigo
        'secondary': RGBColor(99, 102, 241),   # Light indigo
        'accent': RGBColor(16, 185, 129),      # Green
        'text_dark': RGBColor(31, 41, 55),     # Dark gray
        'text_light': RGBColor(255, 255, 255), # White
        'background': RGBColor(249, 250, 251), # Light gray
    }
    
    SLIDE_WIDTH = Inches(13.333)  # 16:9 aspect ratio
    SLIDE_HEIGHT = Inches(7.5)
    
    def __init__(self, branding: Dict[str, str] = None, template_path: str = None):
        """
        Initialize PPT service with optional branding and template.
        
        Args:
            branding: Custom color scheme
            template_path: Path to PPTX template file to use as base
        """
        self.colors = self.DEFAULT_COLORS.copy()
        self.template_path = template_path
        
        if branding:
            if 'primary_color' in branding:
                self.colors['primary'] = self._hex_to_rgb(branding['primary_color'])
            if 'secondary_color' in branding:
                self.colors['secondary'] = self._hex_to_rgb(branding['secondary_color'])
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor."""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    
    def _extract_template_colors(self, prs: Presentation):
        """Extract theme colors from template slide master."""
        try:
            # Try to get colors from the first slide master's theme
            if prs.slide_masters and len(prs.slide_masters) > 0:
                master = prs.slide_masters[0]
                theme = master.slide_master.element
                
                # Log available theme info for debugging
                logger.info(f"Template has {len(prs.slide_masters)} slide masters")
                
                # The theme colors are embedded in the XML, for now we'll just log
                # and keep using template layouts which will inherit the styling
                logger.info("Template colors will be inherited from slide layouts")
        except Exception as e:
            logger.warning(f"Could not extract template colors: {e}")
    
    def _get_best_layout(self, prs: Presentation, slide_type: str):
        """
        Get the best matching layout from template for the slide type.
        Falls back to blank layout if no match found.
        
        Layout matching:
        - cover -> Title Slide (layout 0) or first layout
        - content -> Title and Content (layout 1) or blank
        - agenda -> Title and Content or Section Header
        - etc.
        """
        # Layout type hints mapping slide_type to typical layout indices
        layout_hints = {
            'cover': [0, 5],           # Title Slide, or Title Only
            'agenda': [1, 2],          # Title and Content, Section Header  
            'content': [1, 5, 6],      # Title and Content, Title Only, Blank
            'two_column': [3, 4],      # Two Content, Comparison
            'architecture': [5, 6, 1], # Title Only, Blank, Title and Content
            'timeline': [1, 5],        # Title and Content, Title Only
            'team': [1, 5],            # Title and Content, Title Only
            'pricing': [1, 5],         # Title and Content, Title Only
            'closing': [0, 5],         # Title Slide, Title Only
        }
        
        hints = layout_hints.get(slide_type, [6, 5, 1])  # Default to blank-ish
        
        # Try each hint in order
        for idx in hints:
            if idx < len(prs.slide_layouts):
                return prs.slide_layouts[idx]
        
        # Fallback to last layout (usually blank)
        return prs.slide_layouts[-1] if prs.slide_layouts else prs.slide_layouts[6]
    
    def generate_pptx(
        self,
        slides_data: List[Dict[str, Any]],
        title: str = "Proposal",
        client_name: str = "Client",
        company_name: str = "Company"
    ) -> io.BytesIO:
        """
        Generate a PowerPoint file from slide data.
        
        Args:
            slides_data: List of slide dictionaries with content
            title: Presentation title
            client_name: Client name for cover slide
            company_name: Company/vendor name
            
        Returns:
            BytesIO buffer containing the .pptx file
        """
        # Load from template if available, otherwise create new
        if self.template_path:
            try:
                prs = Presentation(self.template_path)
                logger.info(f"Loaded PPTX template from: {self.template_path}")
                
                # Get existing slide count
                template_slide_count = len(prs.slides)
                logger.info(f"Template has {template_slide_count} existing slides and {len(prs.slide_layouts)} layouts")
                
                # IMPORTANT: Remove all existing slides from template
                # We only want to use the template's layouts/styling, not its content
                # Delete slides in reverse order to avoid index issues
                for i in range(template_slide_count - 1, -1, -1):
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
                
                logger.info(f"Cleared {template_slide_count} template slides, keeping layouts only")
                
                # Extract colors from template if possible
                self._extract_template_colors(prs)
                
            except Exception as e:
                logger.warning(f"Failed to load/clear template, creating new: {e}")
                import traceback
                traceback.print_exc()
                prs = Presentation()
                prs.slide_width = self.SLIDE_WIDTH
                prs.slide_height = self.SLIDE_HEIGHT
        else:
            prs = Presentation()
            prs.slide_width = self.SLIDE_WIDTH
            prs.slide_height = self.SLIDE_HEIGHT
        
        # Track if we're using a template
        self._using_template = self.template_path is not None
        
        for slide_data in slides_data:
            slide_type = slide_data.get('slide_type', 'content')
            
            if slide_type == 'cover':
                self._add_cover_slide(prs, slide_data, title, client_name, company_name)
            elif slide_type == 'agenda':
                self._add_agenda_slide(prs, slide_data)
            elif slide_type == 'two_column':
                self._add_two_column_slide(prs, slide_data)
            elif slide_type == 'architecture':
                self._add_architecture_slide(prs, slide_data)
            elif slide_type == 'timeline':
                self._add_timeline_slide(prs, slide_data)
            elif slide_type == 'team':
                self._add_team_slide(prs, slide_data)
            elif slide_type == 'pricing':
                self._add_pricing_slide(prs, slide_data)
            elif slide_type == 'closing':
                self._add_closing_slide(prs, slide_data, company_name)
            else:
                self._add_content_slide(prs, slide_data)
        
        # Save to buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        logger.info(f"Generated PPT with {len(slides_data)} slides")
        return buffer
    
    def _add_cover_slide(
        self,
        prs: Presentation,
        data: Dict[str, Any],
        title: str,
        client_name: str,
        company_name: str
    ):
        """Add a cover/title slide."""
        # Use template layout if available
        if self._using_template:
            slide_layout = self._get_best_layout(prs, 'cover')
        else:
            slide_layout = prs.slide_layouts[6]  # Blank layout
        
        slide = prs.slides.add_slide(slide_layout)
        
        # Try to use placeholder shapes from template layout
        if self._using_template:
            # Find and populate title placeholder
            for shape in slide.shapes:
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # Title placeholder (usually type 1 or CENTER_TITLE)
                    if ph_type in [1, 3]:  # TITLE or CENTER_TITLE
                        shape.text = data.get('title', title)
                    # Subtitle placeholder (usually type 2)
                    elif ph_type == 2:  # SUBTITLE
                        shape.text = f"Proposal for {client_name}"
            
            # Add company name at bottom
            company_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(6.5),
                Inches(11.5), Inches(0.5)
            )
            company_frame = company_box.text_frame
            company_para = company_frame.paragraphs[0]
            company_para.text = f"Prepared by: {company_name}"
            company_para.font.size = Pt(14)
            company_para.alignment = PP_ALIGN.CENTER
        else:
            # Original custom layout for non-template mode
            # Background shape
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                self.SLIDE_WIDTH, self.SLIDE_HEIGHT
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = self.colors['primary']
            bg_shape.line.fill.background()
            
            # Title
            title_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(2.5),
                Inches(11.5), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = data.get('title', title)
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = self.colors['text_light']
            title_para.alignment = PP_ALIGN.CENTER
            
            # Subtitle
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(4.0),
                Inches(11.5), Inches(0.75)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = f"Proposal for {client_name}"
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = self.colors['text_light']
            subtitle_para.alignment = PP_ALIGN.CENTER
            
            # Company name at bottom
            company_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(6.5),
                Inches(11.5), Inches(0.5)
            )
            company_frame = company_box.text_frame
            company_para = company_frame.paragraphs[0]
            company_para.text = f"Prepared by: {company_name}"
            company_para.font.size = Pt(14)
            company_para.font.color.rgb = self.colors['text_light']
            company_para.alignment = PP_ALIGN.CENTER
    
    def _add_content_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Add a standard content slide with bullets."""
        # Use template layout if available
        if self._using_template:
            slide_layout = self._get_best_layout(prs, 'content')
        else:
            slide_layout = prs.slide_layouts[6]  # Blank layout
        
        slide = prs.slides.add_slide(slide_layout)
        
        bullets = data.get('bullets', [])
        
        if self._using_template:
            # Try to populate placeholders from template
            for shape in slide.shapes:
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # Title placeholder
                    if ph_type == 1:  # TITLE
                        shape.text = data.get('title', 'Content')
                    # Body/content placeholder  
                    elif ph_type == 2 or ph_type == 7:  # BODY or OBJECT
                        if hasattr(shape, 'text_frame'):
                            tf = shape.text_frame
                            for i, bullet in enumerate(bullets[:6]):
                                if i == 0:
                                    para = tf.paragraphs[0]
                                else:
                                    para = tf.add_paragraph()
                                para.text = bullet
                                para.level = 0
        else:
            # Original custom layout
            # Header bar
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                self.SLIDE_WIDTH, Inches(1.2)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = self.colors['primary']
            header.line.fill.background()
            
            # Title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(12), Inches(0.7)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = data.get('title', 'Content')
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = self.colors['text_light']
            
            # Bullets
            if bullets:
                content_box = slide.shapes.add_textbox(
                    Inches(0.75), Inches(1.6),
                    Inches(11.5), Inches(5.5)
                )
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for i, bullet in enumerate(bullets[:6]):  # Max 6 bullets
                    if i == 0:
                        para = content_frame.paragraphs[0]
                    else:
                        para = content_frame.add_paragraph()
                    
                    para.text = f"• {bullet}"
                    para.font.size = Pt(20)
                    para.font.color.rgb = self.colors['text_dark']
                    para.space_before = Pt(12)
                    para.space_after = Pt(8)
    
    def _add_agenda_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Add an agenda slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors['primary']
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Agenda"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_light']
        
        # Agenda items with numbers
        bullets = data.get('bullets', [])
        if bullets:
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.8),
                Inches(11), Inches(5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for i, item in enumerate(bullets[:10], 1):
                if i == 1:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()
                
                para.text = f"{i}. {item}"
                para.font.size = Pt(18)
                para.font.color.rgb = self.colors['text_dark']
                para.space_before = Pt(10)
    
    def _add_two_column_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Add a two-column comparison slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors['primary']
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = data.get('title', 'Comparison')
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_light']
        
        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.6),
            Inches(5.8), Inches(5.5)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        left_items = data.get('left_column', data.get('bullets', [])[:3])
        for i, item in enumerate(left_items):
            if i == 0:
                para = left_frame.paragraphs[0]
            else:
                para = left_frame.add_paragraph()
            para.text = f"• {item}"
            para.font.size = Pt(18)
            para.font.color.rgb = self.colors['text_dark']
            para.space_before = Pt(8)
        
        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.6),
            Inches(5.8), Inches(5.5)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        right_items = data.get('right_column', data.get('bullets', [])[3:6])
        for i, item in enumerate(right_items):
            if i == 0:
                para = right_frame.paragraphs[0]
            else:
                para = right_frame.add_paragraph()
            para.text = f"• {item}"
            para.font.size = Pt(18)
            para.font.color.rgb = self.colors['text_dark']
            para.space_before = Pt(8)
    
    def _add_architecture_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Add an architecture diagram slide with actual rendered diagram."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors['primary']
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = data.get('title', 'Solution Architecture')
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_light']
        
        # Check if we have mermaid code to render
        mermaid_code = data.get('mermaid_code')
        diagram_rendered = False
        
        if mermaid_code:
            try:
                from .mermaid_service import render_mermaid_to_bytes_io
                diagram_buffer = render_mermaid_to_bytes_io(mermaid_code)
                if diagram_buffer:
                    # Add the actual diagram image
                    slide.shapes.add_picture(
                        diagram_buffer,
                        Inches(1), Inches(1.8),
                        width=Inches(11)
                    )
                    diagram_rendered = True
                    logger.info("Architecture diagram rendered successfully in slide")
            except Exception as e:
                logger.error(f"Failed to render diagram in slide: {e}")
        
        if not diagram_rendered:
            # Fallback: Show placeholder
            diagram_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(1.8),
                Inches(11), Inches(4.5)
            )
            diagram_box.fill.solid()
            diagram_box.fill.fore_color.rgb = RGBColor(243, 244, 246)
            diagram_box.line.color.rgb = RGBColor(209, 213, 219)
            
            # Placeholder text
            placeholder_box = slide.shapes.add_textbox(
                Inches(3), Inches(3.5),
                Inches(7), Inches(1)
            )
            placeholder_frame = placeholder_box.text_frame
            placeholder_para = placeholder_frame.paragraphs[0]
            placeholder_para.text = "[Architecture Diagram]"
            placeholder_para.font.size = Pt(24)
            placeholder_para.font.color.rgb = RGBColor(156, 163, 175)
            placeholder_para.alignment = PP_ALIGN.CENTER
        
        # Visual suggestion note (if no diagram was rendered)
        if not diagram_rendered and data.get('visual_suggestion'):
            note_box = slide.shapes.add_textbox(
                Inches(1), Inches(6.5),
                Inches(11), Inches(0.5)
            )
            note_frame = note_box.text_frame
            note_para = note_frame.paragraphs[0]
            note_para.text = f"Suggested: {data['visual_suggestion']}"
            note_para.font.size = Pt(12)
            note_para.font.italic = True
            note_para.font.color.rgb = RGBColor(107, 114, 128)
    
    def _add_timeline_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Add a timeline slide."""
        # Use content slide layout for timeline
        self._add_content_slide(prs, data)
    
    def _add_team_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Add a team structure slide."""
        # Use content slide layout for team
        self._add_content_slide(prs, data)
    
    def _add_pricing_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Add a pricing summary slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Header with accent color
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors['accent']
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = data.get('title', 'Pricing Summary')
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_light']
        
        # Content
        bullets = data.get('bullets', [])
        if bullets:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.6),
                Inches(11.5), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for i, bullet in enumerate(bullets[:6]):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()
                
                para.text = f"• {bullet}"
                para.font.size = Pt(20)
                para.font.color.rgb = self.colors['text_dark']
                para.space_before = Pt(12)
    
    def _add_closing_slide(self, prs: Presentation, data: Dict[str, Any], company_name: str):
        """Add a closing/thank you slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.colors['primary']
        bg_shape.line.fill.background()
        
        # Thank you text
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(2.5),
            Inches(11.5), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = data.get('title', 'Thank You')
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_light']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Q&A text
        qa_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(4.0),
            Inches(11.5), Inches(0.75)
        )
        qa_frame = qa_box.text_frame
        qa_para = qa_frame.paragraphs[0]
        qa_para.text = "Questions & Discussion"
        qa_para.font.size = Pt(24)
        qa_para.font.color.rgb = self.colors['text_light']
        qa_para.alignment = PP_ALIGN.CENTER
        
        # Company name
        company_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(6.0),
            Inches(11.5), Inches(0.5)
        )
        company_frame = company_box.text_frame
        company_para = company_frame.paragraphs[0]
        company_para.text = company_name
        company_para.font.size = Pt(16)
        company_para.font.color.rgb = self.colors['text_light']
        company_para.alignment = PP_ALIGN.CENTER
